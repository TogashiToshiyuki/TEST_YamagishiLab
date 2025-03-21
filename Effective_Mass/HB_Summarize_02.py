#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import argparse
import datetime
import functools
import glob
import math
import os
import time

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from matplotlib.animation import FuncAnimation
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt
from scipy.interpolate import griddata
from scipy.optimize import differential_evolution
from tabulate import tabulate

print = functools.partial(print, flush=True)


def main():
    args, before = get_args()
    print(StandardPhrases.ProgramAbst)
    Operator = input(f"\n{Color.GREEN}Retrieve the operator name.{Color.RESET}\n"
                     f"\t>>> The information entered here will be used as "
                     f"{Color.GREEN}{Color.UNDERLINE}The Name of the Slide Creator{Color.RESET}.\n"
                     f"\tPlease enter your Name.\n"
                     f"\t{Color.GREEN}>>> {Color.RESET}")
    if Operator == "":
        Operator = "ONE"
    else:
        pass
    EM = EffectiveMass(args)

    # Make band information files
    print("\n**********\nMaking band information files...")
    if len(EM.p1Files) != 0 and len(EM.p2Files) != 0:
        DatList_p12 = EM.mkBandInfo_p1p2()
    else:
        DatList_p12 = []
    if len(EM.p3Files) != 0:
        DatList_p3 = EM.mkBandInfo_p3()
    else:
        DatList_p3 = []
    DatList = DatList_p12 + DatList_p3
    DatList.sort()
    print(f"\n\t>>> {Color.GREEN}The band information files are created.{Color.RESET}")

    # Make band structure figures
    print("\n**********\nMaking band structures...")
    Pattern = Constants.Pattern
    FailFiles = []
    MassValues = []
    plotPNGName = ""
    for Dat in DatList:
        print(f"\n>>> Creating band structure figures for {Color.GREEN}{Dat}{Color.RESET}...")
        Params, Fail = EM.getParameters(Dat)
        if Fail:
            FailFiles.append(Dat)
        else:
            if EM.debug:
                EM.displayDef(Params)
            dev, Kcol, Ktrv, Energy_plus_array, Energy_minus_array, Mass, Masses = EM.calcEffMass(Params)
            MassValues.append(EM.mkEffectiveMassLine(Dat, Params, Masses))
            plotPNGName = EM.plotBandDisp(Params["Comment"], dev, Kcol, Ktrv, Energy_plus_array, Energy_minus_array,
                                          Masses, Pattern)
    print(f"\n>>> {Color.GREEN}The band structures were ploted for {len(MassValues)} energy bands.{Color.RESET}\n")

    MassValues.sort()

    parts = plotPNGName.split('-')
    trimmed_parts = parts[:-3]
    MaterName_tNd = '-'.join(trimmed_parts)
    MaterName_tNd = MaterName_tNd.split("_")
    header = ("Entry\tAngle\tlattice in column direction\tlattice in transverse direction\t"
              "Mcol from tensor\tMtrv from tensor\tMcol\tMtrv\tMcol/Mtrv")
    with open(f"./{EM.MaterName}_{MaterName_tNd[1]}_EffMasses.txt", "w") as f:
        f.write(f"*************** Effective Masses of {EM.MaterName} ***************\n")
        f.write(header + "\n")
        for MassValue in MassValues:
            f.write(f"{MassValue.strip()}\n")
        columns = header.strip().split('\t')
        data = []
        for MassValue in MassValues:
            fields = MassValue.strip().split('\t')
            data.append(fields)
        df = pd.DataFrame(data, columns=columns)
        numeric_cols = columns[1:]
        for col in numeric_cols:
            df[col] = pd.to_numeric(df[col], errors='coerce').fillna(df[col])
        print(tabulate(df, headers='keys', tablefmt='pretty', showindex=False))

    print(f"\n\t>>> {Color.GREEN}Effective Masses were evaluated for '{len(MassValues)}' energy bands.{Color.RESET}")

    print("\n**********\nMaking summary plots...")
    EM.Make_Summary_Plots(MaterName_tNd[1])

    print("\n**********\nMaking summary slide...")
    PPTX = SummarySlide(args, Operator, MaterName_tNd[1])
    OUTPUT_FILE_PATH = f"{PPTX.MaterName}_{MaterName_tNd[1]}_summary.pptx"
    print(f"\t>>> Output File: {Color.GREEN}{OUTPUT_FILE_PATH}{Color.RESET}")

    PPTX.mkFrontCover()

    for i in range(2):
        slide = PPTX.add_slide()
        PPTX.add_title(slide)
        shapes = slide.shapes
        shape = shapes.add_textbox(0, Constants.SLIDE_HEIGHT / 30, Constants.SLIDE_WIDTH, Constants.SLIDE_HEIGHT / 10)
        text_frame = shape.text_frame
        text_frame.clear()
        pg = text_frame.paragraphs[0]
        pg.font.size = Pt(36)
        pg.alignment = PP_ALIGN.CENTER
        pg.text = f"Slide {i + 1}"
        pg.font.name = "メイリオ"

    PPTX.add_EvsAngles()
    PPTX.add_TIs()
    PPTX.addStableBandFigures()
    PPTX.add_EMvsAngles()

    # 有効質量が小さいバンド構造のテンプレートを作成
    slide = PPTX.add_slide()
    PPTX.add_title(slide)
    shapes = slide.shapes
    shape = shapes.add_textbox(0, Constants.SLIDE_HEIGHT / 30, Constants.SLIDE_WIDTH, Constants.SLIDE_HEIGHT / 10)
    text_frame = shape.text_frame
    text_frame.clear()
    pg = text_frame.paragraphs[0]
    pg.font.size = Pt(36)
    pg.font.name = "メイリオ"
    pg.alignment = PP_ALIGN.CENTER
    pg.text = "有効質量が小さいバンド構造"

    # pptxファイルを出力する
    PPTX.prs.save(OUTPUT_FILE_PATH)

    after = time.time()
    elapsed_time = after - before
    minutes = int(elapsed_time // 60)
    seconds = int(elapsed_time % 60)
    print(f"\n"
          f"Elapsed Time: {minutes} min {seconds} s\n{Color.GREEN}"
          f"************************* ALL PROCESSES END *************************"
          f"{Color.RESET}\n")


def get_args():
    before = time.time()
    help_text = StandardPhrases.ProgramAbst

    parser = argparse.ArgumentParser(description=help_text, formatter_class=argparse.RawTextHelpFormatter)
    parser.add_argument('MaterName',
                        help='Material name')
    parser.add_argument('--debug', '-d', '-D',
                        help="Start in debug mode",
                        action='store_true')
    parser.add_argument('--gif', '-g', '-G',
                        help="Create GJF File from All File",
                        action='store_true')
    args = parser.parse_args()

    if args.debug:
        print(f"{Color.RED}*************** Caution!!! Debug Started!!! ***************{Color.RESET}\n")

    return args, before


class EffectiveMass:
    def __init__(self, args):
        self.args = args
        self.MaterName = args.MaterName
        self.debug = args.debug
        self.gif = args.gif
        self.messages = []
        self.HelpList = []
        self.p1Files, self.p2Files, self.p3Files, self.AllFiles, self.Angles, self.Tilt_Angle = self.File_Set_Check()
        os.makedirs("./BandInfo", exist_ok=True)
        os.makedirs("./Figures/BandStructures", exist_ok=True)
        os.makedirs("./Figures/EvsAngle", exist_ok=True)
        os.makedirs("./Figures/TIvsAngle", exist_ok=True)
        os.makedirs("./Figures/AllData", exist_ok=True)
        os.makedirs("./Figures/MassvsAngles", exist_ok=True)

    # Displays messages and exits if unexpected behavior is detected.
    def help_check_exit(self):
        """
        Checks for unexpected behaviors and displays accumulated messages.

        This function performs two main tasks:

        1. It displays accumulated messages stored in `self.message`, allowing
           for effective message management within the program.
        2. It checks the `HelpList` attribute for any instances of `True`.
           If `True` is found, indicating an unexpected behavior, the program
           terminates after displaying an error message.

        :returns: None
        :raises SystemExit: Terminates the program if `True` is found in
                            `HelpList`, signaling an abnormal end due to an
                            unexpected behavior.
        """
        self.message_show()
        if True in self.HelpList:
            print(f"{Color.RED}{StandardPhrases.AbnormalEnd}{Color.RESET}")
            exit()
        else:
            pass
        return None

    def message_show(self):
        """
        Show messages.
        :return:
        """
        for message in self.messages:
            print(message)
        self.messages.clear()
        return None

    def File_Set_Check(self):
        """
        Check if the required files exist.
        :return:
        """
        print(f"{Color.GREEN}Checking the required files...{Color.RESET}")
        p1Files, p2Files, p3Files = self.getMinTIFileName()
        AllFiles = self.getAllFileNames()
        self.help_check_exit()
        print(f"{Color.GREEN}All required files are available.{Color.RESET}")

        print("\nTI files:")
        for p1File in p1Files:
            print(f"\t{p1File}")
        for p2File in p2Files:
            print(f"\t{p2File}")
        for p3File in p3Files:
            print(f"\t{p3File}")

        print("\nAll files:")
        for AllFile in AllFiles:
            print(f"\t{AllFile}")

        Angles = []
        for AllFile in AllFiles:
            Angles_temp = self.getAngle(AllFile)
            Angles = Angles + Angles_temp
        Angles = list(set(Angles))
        Angles.sort()
        if self.debug:
            print(f"\nAngles: {Angles}")

        print(f"\nGetting the tilt angles...")
        Tilt_temp = []
        for p1File in p1Files:
            Tilt_temp.append(p1File.split("/")[1].split("_")[2])
        for p2File in p2Files:
            Tilt_temp.append(p2File.split("/")[1].split("_")[2])
        for p3File in p3Files:
            Tilt_temp.append(p3File.split("/")[1].split("_")[2])
        Tilt_temp = list(set(Tilt_temp))
        if len(Tilt_temp) != 1:
            print(f"{Color.RED}\t>>> Error: Tilt angles are not consistent.{Color.RESET}")
            self.HelpList.append(True)
        else:
            print(f"\t>>> Tilt angles: {Color.GREEN}{Tilt_temp[0]}{Color.RESET}")
        return p1Files, p2Files, p3Files, AllFiles, Angles, Tilt_temp[0]

    def getMinTIFileName(self):
        Dirs = os.listdir("./")
        p1Files, p2Files, p3Files = [], [], []
        for Dir in Dirs:
            if self.MaterName in Dir and "_results" in Dir:
                Files_temp = os.listdir(Dir)
                for File in Files_temp:
                    if "3molp1" in File and "min-TIs-" in File:
                        p1Files.append("./" + Dir + "/" + File)
                    elif "3molp2" in File and "min-TIs-" in File:
                        p2Files.append("./" + Dir + "/" + File)
                    elif "3molp3" in File and "min-TIs-" in File:
                        p3Files.append("./" + Dir + "/" + File)
            else:
                pass

        if len(p1Files) == 0:
            print(f"{Color.RED}\t>>> Error: No files found for 3molp1.{Color.RESET}")
            self.HelpList.append(True)
        elif len(p1Files) != 3:
            print(f"{Color.RED}\t>>> Error: 3molp1 files are not complete.{Color.RESET}")
            self.HelpList.append(True)
        else:
            print(f"\t>>> 3molp1 files: {Color.GREEN}All Available.{Color.RESET}")

        if len(p2Files) == 0:
            print(f"{Color.RED}\t>>> Error: No files found for 3molp2.{Color.RESET}")
            self.HelpList.append(True)
        elif len(p2Files) != 3:
            print(f"{Color.RED}\t>>> Error: 3molp2 files are not complete.{Color.RESET}")
            self.HelpList.append(True)
        else:
            print(f"\t>>> 3molp2 files: {Color.GREEN}All Available.{Color.RESET}")

        if len(p3Files) == 0:
            print(f"{Color.RED}\t>>> Error: No files found for 3molp3.{Color.RESET}")
            self.HelpList.append(True)
        elif len(p3Files) != 3:
            print(f"{Color.RED}\t>>> Error: 3molp3 files are not complete.{Color.RESET}")
            self.HelpList.append(True)
        else:
            print(f"\t>>> 3molp3 files: {Color.GREEN}All Available.{Color.RESET}")

        return p1Files, p2Files, p3Files

    def getAllFileNames(self):
        Dirs = os.listdir("./")
        AllFiles = []
        for Dir in Dirs:
            if self.MaterName in Dir and "_results" in Dir and os.path.isdir(Dir):
                Files_temp = os.listdir(Dir)
                for File in Files_temp:
                    if "3mol" in File and "_all.txt" in File:
                        AllFiles.append("./" + Dir + "/" + File)
                    else:
                        pass
            else:
                pass

        if len(AllFiles) == 0:
            print(f"{Color.RED}\t>>> Error: No files found for 3mol_all.txt{Color.RESET}")
            self.HelpList.append(True)
        else:
            print(f"\t>>> 3mol_all.txt files: {Color.GREEN}All Available.{Color.RESET}")

        return AllFiles

    @staticmethod
    def getAngle(AllFile):
        with open(AllFile, "r") as f:
            lines = f.readlines()
        Angles_temp = []
        for line in lines:
            Contents = line.strip().split()
            try:
                Angles_temp.append(float(Contents[1]))
            except (ValueError, IndexError):
                pass
        return Angles_temp

    def mkBandInfo_p1p2(self):

        p1Data_12, p1Data_23, p1Data_31 = [], [], []
        for p1File in self.p1Files:
            if "-12.txt" in p1File:
                p1Data_12 = self.TextFileToData(p1File)
            elif "-23.txt" in p1File:
                p1Data_23 = self.TextFileToData(p1File)
            elif "-31.txt" in p1File:
                p1Data_31 = self.TextFileToData(p1File)

        p2Data_12, p2Data_23, p2Data_31 = [], [], []
        for p2File in self.p2Files:
            if "-12.txt" in p2File:
                p2Data_12 = self.TextFileToData(p2File)
            elif "-23.txt" in p2File:
                p2Data_23 = self.TextFileToData(p2File)
            elif "-31.txt" in p2File:
                p2Data_31 = self.TextFileToData(p2File)

        DatList_temp = []

        for Angle in self.Angles:
            name1, name2 = [], []
            for line in p1Data_12:
                DataList = line.strip().split()
                if float(DataList[1]) == Angle:
                    name1 = DataList[0].split("_")
                    Dcol = float(DataList[2].strip())
                    Dtrv1 = float(DataList[3].strip())
                    T12_HOMO = float(DataList[15].strip())
                    T12_LUMO = float(DataList[13])

            for line in p1Data_23:
                DataList = line.split()
                if float(DataList[1]) == Angle:
                    T23_HOMO = float(DataList[15].strip())
                    T23_LUMO = float(DataList[13].strip())
                else:
                    pass

            for line in p1Data_31:
                DataList = line.split()
                if float(DataList[1]) == Angle:
                    T13_HOMO = float(DataList[15].strip())
                    T13_LUMO = float(DataList[13].strip())
                else:
                    pass

            for line in p2Data_12:
                DataList = line.strip().split()
                if float(DataList[1]) == Angle:
                    name2 = DataList[0].split("_")
                    Dtrv2 = float(DataList[3].strip())
                else:
                    pass

            for line in p2Data_23:
                DataList = line.split()
                if float(DataList[1]) == Angle:
                    T35_HOMO = float(DataList[15].strip())
                    T35_LUMO = float(DataList[13].strip())
                else:
                    pass

            for line in p2Data_31:
                DataList = line.split()
                if float(DataList[1]) == Angle:
                    T34_HOMO = float(DataList[15].strip())
                    T34_LUMO = float(DataList[13].strip())
                else:
                    pass

            # 同じ物質のデータかどうかを確認
            if name1[0] == name2[0]:
                name = name1
            else:
                self.messages.append(f"{Color.RED}"
                                     f"\t>>> Error: Data for the Different Material might exist."
                                     f"{Color.RESET}")
                self.HelpList.append(True)
            self.help_check_exit()

            Dtrv = Dtrv1 + Dtrv2
            title = f"{name[0]}_{self.Tilt_Angle}-B12-{int(Angle)}d"
            print(f"\t>>> {title}-HOMO.dat: ", end="")
            with open(f"./BandInfo/{title}-HOMO.dat", "w") as Fhomo:
                Fhomo.write(f"{title}-HOMO\n")
                Fhomo.write(f"{Constants.n}\n")
                Fhomo.write(f"{Dcol}\n")
                Fhomo.write(f"{Dtrv}\n")
                Fhomo.write("\n")
                Fhomo.write(f"{T12_HOMO}\n")
                Fhomo.write(f"{T13_HOMO}\n")
                Fhomo.write(f"{T23_HOMO}\n")
                Fhomo.write(f"{T34_HOMO}\n")
                Fhomo.write(f"{T35_HOMO}\n")
                Fhomo.write("\n")
            DatList_temp.append(f"./BandInfo/{title}-HOMO.dat")
            print("Complete!")

            print(f"\t>>> {title}-LUMO.dat: ", end="")
            with open(f"./BandInfo/{title}-LUMO.dat", "w") as Flumo:
                Flumo.write(f"{title}-LUMO\n")
                Flumo.write(f"{Constants.n}\n")
                Flumo.write(f"{Dcol}\n")
                Flumo.write(f"{Dtrv}\n")
                Flumo.write("\n")
                Flumo.write(f"{T12_LUMO}\n")
                Flumo.write(f"{T13_LUMO}\n")
                Flumo.write(f"{T23_LUMO}\n")
                Flumo.write(f"{T34_LUMO}\n")
                Flumo.write(f"{T35_LUMO}\n")
                Flumo.write("\n")
            DatList_temp.append(f"./BandInfo/{title}-LUMO.dat")
            print("Complete!")

        return DatList_temp

    def mkBandInfo_p3(self):

        p3Data_12, p3Data_23, p3Data_13 = [], [], []
        for p3 in self.p3Files:
            if "-12.txt" in p3:
                p3Data_12 = self.TextFileToData(p3)
            elif "-23.txt" in p3:
                p3Data_23 = self.TextFileToData(p3)
            elif "-31.txt" in p3:
                p3Data_13 = self.TextFileToData(p3)
            else:
                pass

        DatList_temp = []

        for Angle in self.Angles:
            for line in p3Data_12:
                DataList = line.strip().split()
                if float(DataList[1]) == Angle:
                    name = DataList[0].split("_")
                    Dcol = float(DataList[2].strip())
                    Dtrv = float(DataList[3].strip()) * 2
                    T12_HOMO = float(DataList[15])
                    T12_LUMO = float(DataList[13])
                else:
                    pass

            for line in p3Data_23:
                DataList = line.strip().split()
                if float(DataList[1]) == Angle:
                    T23_HOMO = float(DataList[15])
                    T23_LUMO = float(DataList[13])
                else:
                    pass

            for line in p3Data_13:
                DataList = line.strip().split()
                if float(DataList[1]) == Angle:
                    T13_HOMO = float(DataList[15])
                    T13_LUMO = float(DataList[13])
                else:
                    pass

            title = f"{name[0]}_{self.Tilt_Angle}-B3-{int(Angle)}d"
            print(f"\t>>> {title}-HOMO.dat: ", end="")
            with open(f"./BandInfo/{title}-HOMO.dat", "w") as Fhomo:
                Fhomo.write(f"{title}-HOMO\n")
                Fhomo.write(f"{Constants.n}\n")
                Fhomo.write(f"{Dcol}\n")
                Fhomo.write(f"{Dtrv}\n")
                Fhomo.write("\n")
                Fhomo.write(f"{T12_HOMO}\n")
                Fhomo.write(f"{T13_HOMO}\n")
                Fhomo.write(f"{T23_HOMO}\n")
                Fhomo.write(f"{T13_HOMO}\n")
                Fhomo.write(f"{T23_HOMO}\n")
                Fhomo.write("\n")
            DatList_temp.append(f"./BandInfo/{title}-HOMO.dat")
            print("Complete!")

            print(f"\t>>> {title}-LUMO.dat: ", end="")
            with open(f"./BandInfo/{title}-LUMO.dat", "w") as Flumo:
                Flumo.write(f"{title}-LUMO\n")
                Flumo.write(f"{Constants.n}\n")
                Flumo.write(f"{Dcol}\n")
                Flumo.write(f"{Dtrv}\n")
                Flumo.write("\n")
                Flumo.write(f"{T12_LUMO}\n")
                Flumo.write(f"{T13_LUMO}\n")
                Flumo.write(f"{T23_LUMO}\n")
                Flumo.write(f"{T13_LUMO}\n")
                Flumo.write(f"{T23_LUMO}\n")
                Flumo.write("\n")
            DatList_temp.append(f"./BandInfo/{title}-LUMO.dat")
            print("Complete!")

        return DatList_temp

    @staticmethod
    def TextFileToData(path):
        with open(path, "r") as f:
            lines = f.readlines()
        del lines[0:2]
        return lines

    def getParameters(self, path):
        if not os.path.isfile(path):
            print(f"\t\t{Color.RED}>>> Error: {path} DOES NOT exist in the current directory.{Color.RESET}")
            self.HelpList.append(True)
            self.help_check_exit()
        else:
            pass
        with open(path, "r") as f:
            lines = f.readlines()
        if lines[4].isspace() and lines[10].isspace():
            Comment = lines[0].strip()
            dev = int(lines[1].strip())
            Dcol = round(float(lines[2].strip()), 3)
            Dtrv = round(float(lines[3].strip()), 3)
            TI12 = round(float(lines[5].strip()), 3)
            TI13 = round(float(lines[6].strip()), 3)
            TI23 = round(float(lines[7].strip()), 3)
            TI34 = round(float(lines[8].strip()), 3)
            TI35 = round(float(lines[9].strip()), 3)
            Fail = False
        else:
            print(f"\t\t{Color.RED}>>> Error: {path} is not a valid file.{Color.RESET}")
            Fail = True
            Comment, dev, Dcol, Dtrv, TI12, TI13, TI23, TI34, TI35 = "", 0, 0, 0, 0, 0, 0, 0, 0
        Params = {"Comment": Comment, "dev": dev, "Dcol": Dcol, "Dtrv": Dtrv, "TI12": TI12, "TI13": TI13, "TI23": TI23,
                  "TI34": TI34, "TI35": TI35}
        return Params, Fail

    @staticmethod
    def displayDef(Params):
        Dcol = Params["Dcol"]
        Dtrv = Params["Dtrv"]
        TI12 = Params["TI12"]
        TI13 = Params["TI13"]
        TI23 = Params["TI23"]
        TI34 = Params["TI34"]
        TI35 = Params["TI35"]
        Comment = Params["Comment"]

        print(f"\t*************** Definition of Variables ***************\n"
              f"\n"
              f"\t                       –Dtransv–\n"
              f"\t(Mol1) ––––––––––––––––––––––––––––––––––––– (Mol4) \n"
              f"\t  |     .                                  .   |\n"
              f"\t  |         . TI13               TI34  .       |\n"
              f"\t  |             .                 .            |\n"
              f"\t  |                 .        .                 |   |\n"
              f"\t TI12                 (Mol3)                   |  Dcol\n"
              f"\t  |                 .        .                 |   |\n"
              f"\t  |             .                 .            |\n"
              f"\t  |        .  TI23              TI35   .       |\n"
              f"\t  |    .                                   .   |\n"
              f"\t(Mol2) ––––––––––––––––––––––––––––––––––––– (Mol5) \n"
              f"\n"
              f"\t  Data from '{Comment}'\n"
              f"\t  Distance in col. = {Dcol} Å\n"
              f"\t  Distance in transv. = {Dtrv} Å\n"
              f"\t  Transfer integrals:\n"
              f"\t       TI12 (Mol1-Mol2) = {TI12} meV\n"
              f"\t       TI13 (Mol1-Mol3) = {TI13} meV\n"
              f"\t       TI23 (Mol2-Mol3) = {TI23} meV\n"
              f"\t       TI34 (Mol3-Mol4) = {TI34} meV\n"
              f"\t       TI35 (Mol3-Mol5) = {TI35} meV\n"
              f"\t*******************************************************\n")
        return None

    def calcEffMass(self, Params):
        print(f"\tCalculating effective masses...")
        Comment = Params["Comment"]
        dev = Params["dev"]
        Dcol = Params["Dcol"]
        Dtrv = Params["Dtrv"]
        TI12 = Params["TI12"]
        TI13 = Params["TI13"]
        TI23 = Params["TI23"]
        TI34 = Params["TI34"]
        TI35 = Params["TI35"]

        # k空間の範囲を定義
        Kcol_range = [-math.pi / Dcol, math.pi / Dcol]
        Ktrv_range = [-math.pi / Dtrv, math.pi / Dtrv]

        # エネルギーを最適化する関数を定義
        def energy_func(k):
            Kc, Kt = k
            results = self.calcEnergy(Dcol, Dtrv, TI12, TI13, TI23, TI34, TI35, Kc, Kt)
            if "HOMO" in Comment:
                return -results[0]  # 最大化するためにマイナスを付ける
            elif "LUMO" in Comment:
                return results[1]  # 最小化
            else:
                print("Error: Comment should contain HOMO or LUMO")
                return None

        # 制約条件を設定（kの範囲内に収める）
        bounds = [Kcol_range, Ktrv_range]

        # 最適化を実行
        res = differential_evolution(energy_func, bounds)

        if not res.success:
            print(f"\t{Color.RED}>>> Optimization failed.{Color.RESET}")
            return None, None, None, None, None, None, None

        # 最適化結果からkベクトルを取得
        Kc_be, Kt_be = res.x

        # エネルギー面を可視化
        self.plot_EnergySurface(Dcol, Dtrv, TI12, TI13, TI23, TI34, TI35, Comment, Kc_be, Kt_be)

        # バンド端のエネルギーを計算
        results = self.calcEnergy(Dcol, Dtrv, TI12, TI13, TI23, TI34, TI35, Kc_be, Kt_be)
        if "HOMO" in Comment:
            Energy = results[0]
            d2E_dKc2 = results[2]
            d2E_dKt2 = results[3]
            d2E_dKcKt = results[4]
        elif "LUMO" in Comment:
            Energy = results[1]
            d2E_dKc2 = results[5]
            d2E_dKt2 = results[6]
            d2E_dKcKt = results[7]
        else:
            Energy, d2E_dKc2, d2E_dKt2, d2E_dKcKt = 0, 0, 0, 0

        print(f"\t\t>>> Energy of the band edge: {Energy} meV at ({round(Kc_be, 6)}, {round(Kt_be, 6)})")

        # 有効質量テンソルを構築
        Mass_tensor = np.array([[d2E_dKc2, d2E_dKcKt],
                                [d2E_dKcKt, d2E_dKt2]])

        # 有効質量の計算
        try:
            # テンソルの対角化
            eigenvalues, eigenvectors = np.linalg.eig(Mass_tensor)
            Mass = Constants.h_bar ** 2 / eigenvalues / Constants.ElMass * 1e20
            Masses = [Mass[0], Mass[1], "from tensor"]
            if self.debug:
                print(f"\t\t>>> d2E_dKc2: {d2E_dKc2}")
                print(f"\t\t>>> d2E_dKt2: {d2E_dKt2}")
                print(f"\t\t>>> d2E_dKcKt: {d2E_dKcKt}")
            print(f"\t>>> Mass in {Color.BLUE}Column{Color.RESET}: {Mass[0]} m0 "
                  f"{Color.GREEN}from tensor{Color.RESET}")
            print(f"\t>>> Mass in {Color.BLUE}Transverse{Color.RESET}: {Mass[1]} m0 "
                  f"{Color.GREEN}from tensor{Color.RESET}")
        except Exception as e:
            print(f"{Color.RED}\t>>> Error: {e}{Color.RESET}")
            Mass_col = Constants.h_bar ** 2 / d2E_dKc2 / Constants.ElMass * 1e20
            Mass_trv = Constants.h_bar ** 2 / d2E_dKt2 / Constants.ElMass * 1e20
            if self.debug:
                print(f"\t\t>>> d2E_dKc2: {d2E_dKc2}")
                print(f"\t\t>>> d2E_dKt2: {d2E_dKt2}")
                print(f"\t\t>>> d2E_dKcKt: {d2E_dKcKt}")
            print(f"\t>>> Mass in {Color.BLUE}Column{Color.RESET}: {Mass_col} m0 "
                  f"from {Color.RED}NOT{Color.RESET} tensor")
            print(f"\t>>> Mass in {Color.BLUE}Transverse{Color.RESET}: {Mass_trv} m0 "
                  f"from {Color.RED}NOT{Color.RESET} tensor")
            Mass = np.array([[], []])
            Masses = [Mass_col, Mass_trv, "not from tensor"]

        # バンド図のためのデータを返す
        Kcol_for_plot = round(math.pi / Dcol / dev, 10)
        Ktrv_for_plot = round(math.pi / Dtrv / dev, 10)
        Energy_plus_array, Energy_minus_array = np.zeros((dev + 1, dev + 1)), np.zeros((dev + 1, dev + 1))
        for i in range(dev + 1):
            for j in range(dev + 1):
                Kc = Kcol_for_plot * i
                Kt = Ktrv_for_plot * j
                Energy_plus, Energy_minus, *_ = self.calcEnergy(Dcol, Dtrv, TI12, TI13, TI23, TI34, TI35, Kc, Kt)
                Energy_plus_array[i][j] = Energy_plus
                Energy_minus_array[i][j] = Energy_minus

        return dev, Kcol_for_plot, Ktrv_for_plot, Energy_plus_array, Energy_minus_array, Mass, Masses

    @staticmethod
    def calcEnergy(column, transv, TI12, TI13, TI23, TI34, TI35, Kcol, Ktrv):
        B11 = 2 * TI12 * math.cos(Kcol * column)
        B12R = ((TI13 + TI35) * math.cos(Kcol * (column / 2) + Ktrv * (transv / 2))
                + (TI23 + TI34) * math.cos(Kcol * (column / 2) - Ktrv * (transv / 2)))
        B12I = ((TI35 - TI13) * math.sin(Kcol * (column / 2) + Ktrv * (transv / 2))
                + (TI23 - TI34) * math.sin(Kcol * (column / 2) - Ktrv * (transv / 2)))
        B12 = math.sqrt(B12R ** 2 + B12I ** 2)
        # Avoid division by zero
        epsilon = 1e-20
        B12 = max(B12, epsilon)

        Energy_plus = B11 + B12
        Energy_minus = B11 - B12

        # Kcolで一階微分
        # dB11_dKc = -2 * TI12 * column * math.sin(column * Kcol)
        dB12R_dKc = -1 * (column / 2) * (
                (TI13 + TI35) * math.sin(((column / 2) * Kcol) + ((transv / 2) * Ktrv)) +
                (TI23 + TI34) * math.sin(((column / 2) * Kcol) - ((transv / 2) * Ktrv))
        )
        dB12I_dKc = (column / 2) * (
                (TI35 - TI13) * math.cos(((column / 2) * Kcol) + ((transv / 2) * Ktrv)) +
                (TI23 - TI34) * math.cos(((column / 2) * Kcol) - ((transv / 2) * Ktrv))
        )
        # dB12_dKc = (1 / B12) * (dB12R_dKc * B12R + dB12I_dKc * B12I)
        # dEplus_dKc = dB11_dKc + dB12_dKc
        # dEminus_dKc = dB11_dKc - dB12_dKc

        # Ktrvで一階微分
        # dB11_dKt = 0
        dB12R_dKt = -1 * (transv / 2) * (
                (TI13 + TI35) * math.sin(((column / 2) * Kcol) + ((transv / 2) * Ktrv)) -
                (TI23 + TI34) * math.sin(((column / 2) * Kcol) - ((transv / 2) * Ktrv))
        )
        dB12I_dKt = (transv / 2) * (
                (TI35 - TI13) * math.cos(((column / 2) * Kcol) + ((transv / 2) * Ktrv)) -
                (TI23 - TI34) * math.cos(((column / 2) * Kcol) - ((transv / 2) * Ktrv))
        )
        # dB12_dKt = (1 / B12) * (dB12R_dKt * B12R + dB12I_dKt * B12I)
        # dEplus_dKt = dB11_dKt + dB12_dKt
        # dEminus_dKt = dB11_dKt - dB12_dKt

        # Kcolで二階微分
        d2B11_dKc2 = -2 * TI12 * (column ** 2) * math.cos(column * Kcol)
        d2B12R_dKc2 = -(column / 2) ** 2 * (
                (TI13 + TI35) * math.cos((column / 2) * Kcol + (transv / 2) * Ktrv) +
                (TI23 + TI34) * math.cos((column / 2) * Kcol - (transv / 2) * Ktrv)
        )
        d2B12I_dKc2 = -(column / 2) ** 2 * (
                (TI35 - TI13) * math.sin((column / 2) * Kcol + (transv / 2) * Ktrv) +
                (TI23 - TI34) * math.sin((column / 2) * Kcol - (transv / 2) * Ktrv)
        )
        d2B12_dKc2 = (1 / B12) * (
                (dB12R_dKc ** 2) + (dB12I_dKc ** 2) + B12R * d2B12R_dKc2 + B12I * d2B12I_dKc2
        ) - ((1 / B12) ** 3) * (
                             (B12R * dB12R_dKc + B12I * dB12I_dKc) ** 2
                     )
        d2Eplus_dKc2 = d2B11_dKc2 + d2B12_dKc2
        d2Eminus_dKc2 = d2B11_dKc2 - d2B12_dKc2

        # Ktrvで二階微分
        d2B11_dKt2 = 0
        d2B12R_dKt2 = -(transv / 2) ** 2 * (
                (TI13 + TI35) * math.cos((column / 2) * Kcol + (transv / 2) * Ktrv) +
                (TI23 + TI34) * math.cos((column / 2) * Kcol - (transv / 2) * Ktrv)
        )
        d2B12I_dKt2 = -(transv / 2) ** 2 * (
                (TI35 - TI13) * math.sin((column / 2) * Kcol + (transv / 2) * Ktrv) +
                (TI23 - TI34) * math.sin((column / 2) * Kcol - (transv / 2) * Ktrv)
        )
        d2B12_dKt2 = (1 / B12) * ((dB12R_dKt ** 2) + (dB12I_dKt ** 2) + B12R * d2B12R_dKt2 + B12I * d2B12I_dKt2) - (
                (1 / B12) ** 3) * (
                             (B12R * dB12R_dKt + B12I * dB12I_dKt) ** 2
                     )
        d2Eplus_dKt2 = d2B11_dKt2 + d2B12_dKt2
        d2Eminus_dKt2 = d2B11_dKt2 - d2B12_dKt2

        # KcolとKtrvでの混合微分
        d2B11_dKcKt = 0
        d2B12R_dKcdKt = -1 * (column / 2) * (transv / 2) * (
                (TI13 + TI35) * math.cos((column / 2) * Kcol + (transv / 2) * Ktrv) -
                (TI23 + TI34) * math.cos((column / 2) * Kcol - (transv / 2) * Ktrv)
        )
        d2B12I_dKcdKt = -1 * (column / 2) * (transv / 2) * (
                (TI35 - TI13) * math.sin((column / 2) * Kcol + (transv / 2) * Ktrv) -
                (TI23 - TI34) * math.sin((column / 2) * Kcol - (transv / 2) * Ktrv)
        )
        d2B12_dKcKt = (1 / B12) * (
                dB12R_dKc * dB12R_dKt + B12R * d2B12R_dKcdKt + dB12I_dKc * dB12I_dKt + B12I * d2B12I_dKcdKt) - (
                              ((1 / B12) ** 3) * (B12R * dB12R_dKc + B12I * dB12I_dKc) *
                              (B12R * dB12R_dKt + B12I * dB12I_dKt)
                      )
        dEplus_dKcKt = d2B11_dKcKt + d2B12_dKcKt
        dEminus_dKcKt = d2B11_dKcKt - d2B12_dKcKt

        return (Energy_plus, Energy_minus,
                d2Eplus_dKc2, d2Eplus_dKt2, dEplus_dKcKt, d2Eminus_dKc2, d2Eminus_dKt2, dEminus_dKcKt)

    def plot_EnergySurface(self, Dcol, Dtrv, TI12, TI13, TI23, TI34, TI35, Comment, Kc_be, Kt_be):
        if self.debug:
            # k空間を細かく分割
            Kc_values = np.linspace(-math.pi / Dcol, math.pi / Dcol, 200)
            Kt_values = np.linspace(-math.pi / Dtrv, math.pi / Dtrv, 200)
            Kc_grid, Kt_grid = np.meshgrid(Kc_values, Kt_values)
            Energy_grid = np.zeros_like(Kc_grid)

            for i in range(len(Kc_values)):
                for j in range(len(Kt_values)):
                    Kc = Kc_values[i]
                    Kt = Kt_values[j]
                    results = self.calcEnergy(Dcol, Dtrv, TI12, TI13, TI23, TI34, TI35, Kc, Kt)
                    if "HOMO" in Comment:
                        Energy_grid[j, i] = results[0]
                    elif "LUMO" in Comment:
                        Energy_grid[j, i] = results[1]

            plt.figure()
            plt.contourf(Kc_grid, Kt_grid, Energy_grid, levels=50, cmap='viridis')
            plt.colorbar(label='Energy (meV)')
            plt.xlabel('Kc')
            plt.ylabel('Kt')
            plt.title(f'Energy Surface for {Comment}')

            # 極値点をプロット
            plt.plot(Kc_be, Kt_be, '.', markersize=10, label='Band Edge', color='red')
            plt.legend()
            plt.show()
            plt.close()
        else:
            pass
        return None

    @staticmethod
    def mkEffectiveMassLine(Dat, Params, Mass_array):
        Direction_column = Params["Dcol"]
        Direction_transv = Params["Dtrv"]
        Comment = Params["Comment"]
        Angle = Comment.split("-")[-2].split("d")[0]

        if Mass_array[2] == "from tensor":
            line = (f"{Dat}\t{Angle}\t{Direction_column}\t{Direction_transv}\t"
                    f"{Mass_array[0]}\t{Mass_array[1]}\t-\t-\t{Mass_array[0] / Mass_array[1]}\n")
        else:
            line = (f"{Dat}\t{Angle}\t{Direction_column}\t{Direction_transv}\t"
                    f"-\t-\t{Mass_array[0]}\t{Mass_array[1]}\t{Mass_array[0] / Mass_array[1]}\n")
        return line

    @staticmethod
    def EnergyPlotParts(dev, E_plus_array, E_minus_array, xUnitLen, D_col, D_trv):
        x_dev_temp = np.full(dev + 1, xUnitLen)
        x_dev_temp[0] = 0
        y1_temp, y2_temp = [], []
        for i in range(dev + 1):
            y1_temp.append(E_plus_array[int(D_col[i]), int(D_trv[i])])
            y2_temp.append(E_minus_array[int(D_col[i]), int(D_trv[i])])
        return x_dev_temp, y1_temp, y2_temp

    def plotBandDisp(self, Comment, dev, Kcol, Ktrv, Energy_plus_array, Energy_minus_array, Masses, Pattern):
        # Make plot for Energy dispersion
        ZeroTo = np.arange(0, dev + 1, 1)
        ToZero = np.flipud(ZeroTo)
        Zero = np.zeros(dev + 1)
        Top = np.full(dev + 1, dev)

        if "1" in Pattern:
            xdev = []
            y1 = []
            y2 = []
            xUnitLen = Ktrv
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, Top, ToZero)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            xUnitLen = Kcol
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, ToZero, Zero)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            xUnitLen = Ktrv
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, Zero, ZeroTo)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            xUnitLen = Kcol
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, ZeroTo, Top)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            xUnitLen = math.sqrt(Kcol ** 2 + Ktrv ** 2)
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, ToZero, ToZero)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            x, val = [], 0
            for i in range(len(xdev)):
                if i == 0:
                    val = 0 + xdev[i]
                    x.append(val)
                else:
                    val = val + xdev[i]
                    x.append(val)

            point1 = x[0]  # C point: pi/a,pi/b
            point2 = x[dev + 1]  # X point: pi/a,0
            point3 = x[2 * (dev + 1)]  # Gamma point: 0.0, 0.0
            point4 = x[3 * (dev + 1)]  # Y point: 0,pi/b
            point5 = x[4 * (dev + 1)]  # C point: pi/a, pi/b
            point6 = x[-1]  # Gamma point: 0.0, 0.0
            xtickvals = [point1, point2, point3, point4, point5, point6]
            xticktexts = ["C", "X", "gamma", "Y", "C", "gamma"]

        elif "2" in Pattern:
            xdev = []
            y1 = []
            y2 = []

            xUnitLen = Kcol
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, ToZero, Zero)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            xUnitLen = Ktrv
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, Zero, ZeroTo)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            xUnitLen = Kcol
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, ZeroTo, Top)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            xUnitLen = Ktrv
            x_dev_temp, y1_temp, y2_temp = self.EnergyPlotParts(dev, Energy_plus_array, Energy_minus_array,
                                                                xUnitLen, Top, ToZero)
            xdev.extend(x_dev_temp)
            y1.extend(y1_temp)
            y2.extend(y2_temp)

            x, val = [], 0
            for i in range(len(xdev)):
                if i == 0:
                    val = 0 + xdev[i]
                    x.append(val)
                else:
                    val = val + xdev[i]
                    x.append(val)

            # print(len(x))
            point1 = x[0]  # X point: pi/a,0
            point2 = x[dev + 1]  # Gamma point: 0.0, 0.0
            point3 = x[2 * (dev + 1)]  # Y point: 0,pi/b
            point4 = x[3 * (dev + 1)]  # C point: pi/a, pi/b
            point5 = x[-1]  # X point: pi/a,0
            xtickvals = [point1, point2, point3, point4, point5]
            xticktexts = ["X", "gamma", "Y", "C", "X"]

        else:
            print(f"\t{Color.RED}>>> Error: Pattern should be 1 or 2.{Color.RESET}")
            self.HelpList.append(True)
            self.help_check_exit()
            exit()

        plt.figure()
        plt.plot(x, y1, color="r")
        plt.plot(x, y2, color="r")
        for xtickval in xtickvals:
            plt.axvline(xtickval, color="k", linewidth=0.75)
        plt.xticks(xtickvals, xticktexts)
        plt.tick_params(axis="both", direction="in", labelsize=14)
        plt.xlim(xtickvals[0], xtickvals[-1])
        plt.ylabel("E (meV)", fontsize=18)
        plt.title(f"{Comment}: mc={round(Masses[0], 2)}m0, mt={round(Masses[1], 2)}m0 '{Masses[2]}'", fontsize=10)
        plt.tight_layout()
        if self.debug:
            plt.show()
            pass
        else:
            pass

        plotPNGname = Comment[:Comment.find(",")].strip()
        plotPNGname = plotPNGname[:-2]
        plt.savefig(f"./Figures/BandStructures/{plotPNGname}-{Pattern.strip()}.png", format="png",
                    dpi=500, bbox_inches="tight")
        plt.close()

        with open(f"./BandInfo/{plotPNGname}-{Pattern.strip()}.txt", "w") as f:
            for i in range(len(x)):
                f.write(f"{x[i]}\t{y1[i]}\t{y2[i]}\n")

        return plotPNGname

    def List_files(self, Tilt_Angle):
        ContentsList = os.listdir("./")
        FolderList = []
        for Content in ContentsList:
            if os.path.isdir(Content) and self.MaterName in Content and "_results" in Content and "3mol" in Content:
                FolderList.append(Content)
            else:
                pass

        FileList = []
        for Folder in FolderList:
            ContentsList = os.listdir(f"./{Folder}")
            for Content in ContentsList:
                if os.path.isfile(f"./{Folder}/{Content}") and "_all.txt" in Content:
                    FileList.append(f"./{Folder}/{Content}")
                elif os.path.isfile(f"./{Folder}/{Content}"):
                    if "_min-TIs-12.txt" in Content and "3mol" in Content:
                        FileList.append(f"./{Folder}/{Content}")
                    elif "_min-TIs-23.txt" in Content and "3mol" in Content:
                        FileList.append(f"./{Folder}/{Content}")
                    elif "_min-TIs-31.txt" in Content and "3mol" in Content:
                        FileList.append(f"./{Folder}/{Content}")
                    else:
                        pass

        if os.path.isfile(f"./{self.MaterName}_{Tilt_Angle}_EffMasses.txt"):
            FileList.append(f"./{self.MaterName}_{Tilt_Angle}_EffMasses.txt")
        else:
            pass

        return FileList

    def Make_Summary_Plots(self, Tilt_Angle):
        FileList = self.List_files(Tilt_Angle)
        FileList.sort()
        for File in FileList:
            print(f"\t>>> Making plots from {File}: ", end="")
            if "min-TIs" in File:
                Data = self.getMinData(File)
                self.plot_EvsAngle(Data)
                self.plot_TIvsAngle(Data)
            elif "_all.txt" in File:
                with open(File, "r") as f:
                    lines = f.readlines()
                del lines[0:2]

                Degs = []
                for line in lines:
                    contents = line.strip().split()
                    if "***" in contents[0]:
                        pass
                    else:
                        try:
                            Degs.append(round(float(contents[1]), 1))
                        except (ValueError, IndexError):
                            pass
                Degs = list(set(Degs))
                Degs.sort()
                for Deg in Degs:
                    Data = self.getAllData(File, Deg)
                    self.plot_AllData(Data)
            elif "EffMasses" in File:
                Structures = ["-B12-", "-B3-"]
                HOMO_LUMOs = ["HOMO", "LUMO"]
                for Structure in Structures:
                    for HOMO_LUMO in HOMO_LUMOs:
                        Data = self.getEffMassData(File, Structure, HOMO_LUMO)
                        self.plot_EffMass(Data, Structure, HOMO_LUMO)
                        if "B12" in Structure:
                            minData1 = self.getMinData(f"./{self.MaterName}_3molp1_{Tilt_Angle}_results/"
                                                       f"{self.MaterName}_3molp1_{Tilt_Angle}_min-TIs-12.txt")
                            minData2 = self.getMinData(f"./{self.MaterName}_3molp2_{Tilt_Angle}_results/"
                                                       f"{self.MaterName}_3molp2_{Tilt_Angle}_min-TIs-12.txt")
                            minEnergy = np.array(minData1["Energy"]) + np.array(minData2["Energy"])
                            minEnergy = list(minEnergy)
                            minAngle = minData1["Angle"]
                            self.plot_EMassvsAngle(Data, Structure, HOMO_LUMO, minAngle, minEnergy)
                        elif "B3" in Structure:
                            minData3 = self.getMinData(f"./{self.MaterName}_3molp3_{Tilt_Angle}_results/"
                                                       f"{self.MaterName}_3molp3_{Tilt_Angle}_min-TIs-12.txt")
                            minEnergy = np.array(minData3["Energy"]) * 2
                            minEnergy = list(minEnergy)
                            minAngle = minData3["Angle"]
                            self.plot_EMassvsAngle(Data, Structure, HOMO_LUMO, minAngle, minEnergy)

            print(f"{Color.GREEN}Complete!{Color.RESET}")

    @staticmethod
    def getMinData(FileName):
        Angle, Dcol, Dtrv, Energy, TI_HOMO, TI_LUMO = [], [], [], [], [], []
        Data = {}
        name = FileName[FileName.rfind("/") + 1:FileName.rfind("_")]
        if "-12.txt" in FileName:
            name = f"{name}-12"
        elif "-23.txt" in FileName:
            name = f"{name}-23"
        elif "-31.txt" in FileName:
            name = f"{name}-31"
        else:
            pass

        with open(FileName, "r") as f:
            lines = f.readlines()
            del lines[0:2]

        for line in lines:
            contents = line.strip().split()
            try:
                Angle.append(round(float(contents[1]), 1))
                Dcol.append(round(float(contents[2]), 2))
                Dtrv.append(round(float(contents[3]), 3))
                Energy.append(round(float(contents[4]), 5))
                TI_LUMO.append(round(float(contents[13]), 3))
                TI_HOMO.append(round(float(contents[15]), 3))
            except ValueError:
                pass

        Data["Angle"] = Angle
        Data["Dcol"] = Dcol
        Data["Dtrv"] = Dtrv
        Data["Energy"] = Energy
        Data["TI_HOMO"] = TI_HOMO
        Data["TI_LUMO"] = TI_LUMO
        Data["name"] = name
        return Data

    @staticmethod
    def plot_EvsAngle(Data):
        name = Data["name"]
        Angle = Data["Angle"]
        Dcol = Data["Dcol"]
        Dtrv = Data["Dtrv"]
        Energy = Data["Energy"]
        title = f"{name}: Distance, Energy vs Angle"

        fig = plt.figure(figsize=(6, 9))
        ax1 = fig.add_subplot(111)
        ax1.set_title(title, fontsize=16, loc="center", y=1.03)
        ax1.set_xlabel("Angle (deg.)", fontsize=20)
        ax1.set_ylabel("Distance (angstrom)", fontsize=20)
        ax1.set_xlim(100, 0)
        ax1.scatter(Angle, Dcol, label="Dcol", color="b", marker="^", s=100)
        ax1.scatter(Angle, Dtrv, label="Dtrv", color="g", marker="v", s=100)

        ax2 = ax1.twinx()
        ax2.set_ylabel("Energy (a.u.)", fontsize=20)
        ax2.scatter(Angle, Energy, label="Energy", color="r", marker="o", s=100)

        ax1.tick_params(axis="both", direction="in", labelsize=18)
        ax2.tick_params(axis="y", labelcolor="r", colors="r", direction="in", labelsize=18)
        ax1.legend(loc="upper right", bbox_to_anchor=(1.2, 1.05))
        ax2.legend(loc="upper right", bbox_to_anchor=(1.2, 1.1))

        fig.savefig(f"./Figures/EvsAngle/{name}_EvsAngle.png", format="png", dpi=300, bbox_inches="tight")

        plt.close()
        return None

    @staticmethod
    def plot_TIvsAngle(Data):
        name = Data["name"]
        Angle = Data["Angle"]
        TI_HOMO = Data["TI_HOMO"]
        TI_LUMO = Data["TI_LUMO"]
        title = f"{name}: Transfer integral vs Angle"
        Color_homo = "r"
        Color_lumo = "b"

        plt.figure(figsize=(16, 9))
        plt.scatter(Angle, TI_HOMO, c=Color_homo, marker="o", s=100, label="HOMO")
        plt.scatter(Angle, TI_LUMO, c=Color_lumo, marker="s", s=100, label="LUMO")
        plt.tick_params(axis="both", direction="in", labelsize=20)
        plt.xlabel("Angle (deg.)", fontsize=24)
        plt.ylabel("Transfer Integral (meV)", fontsize=24)
        plt.axhline(0, 100, 0, c="k", linestyle="solid", linewidth=0.75)
        plt.xlim(100, 0)
        plt.title(title, fontsize=24, loc="center", y=1.03)
        plt.legend(fontsize=20)
        plt.savefig(f"./Figures/TIvsAngle/{name}_TIs.png", format="png", dpi=400)
        plt.close()
        return

    @staticmethod
    def getAllData(filename, Deg):
        x = []
        y = []
        z = []
        Data = {}
        with open(filename, "r") as file:
            lines = file.readlines()
        del lines[0:2]

        for line in lines:
            contents = line.strip().split()
            if "***" in contents[0]:
                pass
            elif round(float(contents[1]), 1) == round(float(Deg), 1):
                x.append(round(float(contents[2]), 2))
                y.append(round(float(contents[3]), 2))
                z.append(round(float(contents[4]), 6))
            else:
                pass

        Data["Dcol"] = x
        Data["Dtrv"] = y
        Data["Energy"] = z
        Data["name"] = f"{filename[filename.rfind('/') + 1:filename.rfind('.txt')]}-{int(Deg)}"
        return Data

    def plot_AllData(self, Data):
        x = np.array(Data["Dcol"])
        y = np.array(Data["Dtrv"])
        z = np.array(Data["Energy"])
        name = Data["name"]

        if self.gif:
            fig = plt.figure(figsize=(6, 6))
            ax = fig.add_subplot(111, projection="3d")
            ax.set_title(name, size=20)
            ax.set_xlabel(r"$D_{\mathrm{col}}$", fontsize=10)
            ax.set_ylabel(r"$D_{\mathrm{trv}}$", fontsize=10)
            ax.set_zlabel("Energy", fontsize=10)
            ax.tick_params(axis="both", direction="in", labelsize=8)

            # Initialize the plot
            def init():
                ax.scatter(x, y, z, c="r", marker="o")
                return ax,

            # Function to rotate the plot
            def rotate(angle):
                ax.view_init(azim=angle)
                return ax,

            # Create the animation
            ani = FuncAnimation(
                fig,
                rotate,
                frames=np.arange(0, 360, 5),
                init_func=init,
                interval=100,
                blit=False
            )

            # Save the animation as a GIF file
            ani.save(f"./Figures/AllData/{name}.gif", writer="pillow")
        else:
            # グリッドデータを作成
            xi = np.linspace(x.min(), x.max(), 100)
            yi = np.linspace(y.min(), y.max(), 100)
            xi, yi = np.meshgrid(xi, yi)
            zi = griddata((x, y), z, (xi, yi), method='linear')

            # プロット作成
            plt.figure(figsize=(6, 6))
            plt.title(name, size=20)
            plt.xlabel(r"$D_{\mathrm{col}}$", fontsize=10)
            plt.ylabel(r"$D_{\mathrm{trv}}$", fontsize=10)
            plt.tick_params(axis="both", direction="in", labelsize=8)

            # 等高線プロット
            contour = plt.contourf(xi, yi, zi, levels=100, cmap='hot_r')
            plt.colorbar(contour, label='Energy')

            # 元のデータポイントをプロット（オプション）
            plt.scatter(x, y, c='red', s=10, label='Data Points')
            plt.legend()

            # 画像を保存
            plt.savefig(f"./Figures/AllData/{name}_2D.png", format="png", dpi=300, bbox_inches="tight")
            plt.close()

            fig = plt.figure(figsize=(6, 6))
            ax = fig.add_subplot(111, projection="3d")
            ax.set_title(name, size=20)
            ax.set_xlabel(r"$D_{\mathrm{col}}$", fontsize=10)
            ax.set_ylabel(r"$D_{\mathrm{trv}}$", fontsize=10)
            ax.set_zlabel("Energy", fontsize=10)
            ax.tick_params(axis="both", direction="in", labelsize=8)
            ax.scatter(x, y, z, c="r", marker="o")
            fig.savefig(f"./Figures/AllData/{name}.png", format="png", dpi=300, bbox_inches="tight")
            plt.close()

        plt.close()

    @staticmethod
    def getEffMassData(filename, Structure, HorL):
        with open(filename, "r") as file:
            lines = file.readlines()
        del lines[0:2]

        Data = {}
        Angle = []
        Angle2 = []
        Mcol = []
        Mcol2 = []
        Mtrv = []
        Mtrv2 = []
        ratio = []
        rAngle = []
        name = filename[filename.rfind("/") + 1:filename.rfind("_")]
        for line in lines:
            contents = line.strip().split()
            if Structure in contents[0] and HorL in contents[0]:
                if contents[6] == "-" and contents[7] == "-":
                    Angle.append(round(float(contents[1]), 1))
                    Mcol.append(round(float(contents[4]), 3))
                    Mtrv.append(round(float(contents[5]), 3))
                    rAngle.append(round(float(contents[1]), 1))
                    ratio.append(round(float(contents[8]), 7))
                elif contents[4] == "-" and contents[5] == "-":
                    Angle2.append(round(float(contents[1]), 1))
                    Mcol2.append(round(float(contents[6]), 3))
                    Mtrv2.append(round(float(contents[7]), 3))
                    rAngle.append(round(float(contents[1]), 1))
                    ratio.append(round(float(contents[8]), 7))

        Data["Angle"] = Angle
        Data["Mcol"] = Mcol
        Data["Mtrv"] = Mtrv
        Data["ratio"] = ratio
        Data["Angle2"] = Angle2
        Data["Mcol2"] = Mcol2
        Data["Mtrv2"] = Mtrv2
        Data["rAngle"] = rAngle
        Data["name"] = name

        return Data

    @staticmethod
    def plot_EffMass(Data, Structure, HorL):
        name = Data["name"]
        Angle = Data["Angle"]
        Mcol = Data["Mcol"]
        Mtrv = Data["Mtrv"]
        ratio = Data["ratio"]
        Angle2 = Data["Angle2"]
        Mcol2 = Data["Mcol2"]
        Mtrv2 = Data["Mtrv2"]
        rAngle = Data["rAngle"]
        title = f"{name}{Structure}{HorL}: Effective Mass vs Angle"

        fig = plt.figure(figsize=(16, 9))
        ax1 = fig.add_subplot(111)
        ax1.set_title(title, fontsize=16, loc="center", y=1.03)
        ax1.set_xlabel("Angle (deg.)", fontsize=20)
        ax1.set_ylabel(r"$m_{\mathrm{eff}} (m_{\mathrm{0}})$", fontsize=20)
        ax1.set_xlim(100, 0)
        if HorL == "HOMO":
            ax1.set_ylim(-5, 0)
        if HorL == "LUMO":
            ax1.set_ylim(5, 0)
        ax1.scatter(Angle, Mcol, label="Mcol", color="b", marker="^", s=100)
        ax1.scatter(Angle, Mtrv, label="Mtrv", color="g", marker="v", s=100)
        ax1.scatter(Angle2, Mcol2, label="Mcol (not tensor)", color="b", marker="s", s=20)
        ax1.scatter(Angle2, Mtrv2, label="Mtrv (not tensor)", color="g", marker="s", s=20)

        ax2 = ax1.twinx()
        ax2.set_ylim(0, 2)
        ax2.set_ylabel(r"$m_{\mathrm{col}}/m_{\mathrm{trv}}$", fontsize=20)
        ax2.scatter(rAngle, ratio, label="ratio", color="r", marker="o", s=100)

        ax1.tick_params(axis="both", direction="in", labelsize=18)
        ax2.tick_params(axis="y", labelcolor="r", colors="r", direction="in", labelsize=18)
        ax1.legend(loc="upper right", bbox_to_anchor=(1.2, 1.05))
        ax2.legend(loc="upper right", bbox_to_anchor=(1.2, 1.1))

        fig.savefig(f"./Figures/MassvsAngles/{name}{Structure}{HorL}_MassvsAngle.png",
                    format="png", dpi=300, bbox_inches="tight")
        plt.close()
        return

    @staticmethod
    def plot_EMassvsAngle(Data, Structure, HorL, minAngle, minEnergy):
        name = Data["name"]
        Angle = Data["Angle"]
        Mcol = Data["Mcol"]
        Mtrv = Data["Mtrv"]
        Angle2 = Data["Angle2"]
        Mcol2 = Data["Mcol2"]
        Mtrv2 = Data["Mtrv2"]
        title = f"{name}{Structure}{HorL}: Effective Mass, Energy vs Angle"

        fig = plt.figure(figsize=(16, 9))
        ax1 = fig.add_subplot(111)
        ax1.set_title(title, fontsize=16, loc="center", y=1.03)
        ax1.set_xlabel("Angle (deg.)", fontsize=20)
        ax1.set_ylabel(r"$m_{\mathrm{eff}} (m_{\mathrm{0}})$", fontsize=20)
        ax1.set_xlim(100, 0)
        if HorL == "HOMO":
            ax1.set_ylim(-5, 0)
        if HorL == "LUMO":
            ax1.set_ylim(5, 0)
        ax1.scatter(Angle, Mcol, label="Mcol", color="b", marker="^", s=100)
        ax1.scatter(Angle, Mtrv, label="Mtrv", color="g", marker="v", s=100)
        ax1.scatter(Angle2, Mcol2, label="Mcol (not tensor)", color="b", marker="s", s=20)
        ax1.scatter(Angle2, Mtrv2, label="Mtrv (not tensor)", color="g", marker="s", s=20)

        ax2 = ax1.twinx()
        if Structure == "-B12-":
            ax2.set_ylabel(r"$Energy, p1+p2 (a.u.)$", fontsize=20)
        if Structure == "-B3-":
            ax2.set_ylabel(r"$Energy, 2 x p3 (a.u.)$", fontsize=20)
        ax2.scatter(minAngle, minEnergy, label="Energy", color="r", marker="o", s=100)

        ax1.tick_params(axis="both", direction="in", labelsize=18)
        ax2.tick_params(axis="y", labelcolor="r", colors="r", direction="in", labelsize=18)
        ax1.legend(loc="upper right", bbox_to_anchor=(1.2, 1.05))
        ax2.legend(loc="upper right", bbox_to_anchor=(1.2, 1.1))

        fig.savefig(f"./Figures/{name}{Structure}{HorL}_EMvsAngle.png", format="png", dpi=300, bbox_inches="tight")
        plt.close()
        return


class SummarySlide:
    def __init__(self, args, Operator, Tilt_Angle):
        self.args = args
        self.MaterName = args.MaterName
        self.debug = args.debug
        self.messages = []
        self.HelpList = []
        self.prs = Presentation()
        self.prs.slide_width = Constants.SLIDE_WIDTH
        self.prs.slide_height = Constants.SLIDE_HEIGHT
        self.Operator = Operator
        self.IMG_DIR = "./Figures"
        self.Tilt_Angle = Tilt_Angle

    # Displays messages and exits if unexpected behavior is detected.
    def help_check_exit(self):
        """
        Checks for unexpected behaviors and displays accumulated messages.

        This function performs two main tasks:

        1. It displays accumulated messages stored in `self.message`, allowing
           for effective message management within the program.
        2. It checks the `HelpList` attribute for any instances of `True`.
           If `True` is found, indicating an unexpected behavior, the program
           terminates after displaying an error message.

        :returns: None
        :raises SystemExit: Terminates the program if `True` is found in
                            `HelpList`, signaling an abnormal end due to an
                            unexpected behavior.
        """
        self.message_show()
        if True in self.HelpList:
            print(f"{Color.RED}{StandardPhrases.AbnormalEnd}{Color.RESET}")
            exit()
        else:
            pass
        return None

    def message_show(self):
        """
        Show messages.
        :return:
        """
        for message in self.messages:
            print(message)
        self.messages.clear()
        return None

    def mkFrontCover(self):
        """
        表紙のスライドを作成する
        :return:
        """
        slide = self.add_slide()

        self.add_title(slide)

        # テキストボックスの追加(for title)
        shape = slide.shapes.add_textbox(0, Constants.IMG_CENTER_Y / 2 - 800000,
                                         Constants.SLIDE_WIDTH, Constants.IMG_CENTER_Y)
        text = f"集合体構造最適化-{self.MaterName}_{self.Tilt_Angle}"
        text_frame = shape.text_frame
        text_frame.clear()
        # 垂直方向の中央揃えにする
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # パラグラフの追加
        pg = text_frame.paragraphs[0]
        # フォントサイズを大きくする
        pg.font.size = Pt(36)
        # 左右中央揃えにする
        pg.alignment = PP_ALIGN.CENTER
        pg.font.bold = False
        # フォントをメイリオにする
        pg.font.name = "メイリオ"
        pg.text = text

        print(f"\t>>> Title Name: {text}")

        # テキストボックスの追加(for 作成者、作成日時)
        shape = slide.shapes.add_textbox(0, Constants.IMG_CENTER_Y, Constants.SLIDE_WIDTH, Constants.IMG_CENTER_Y)
        Current_Time = datetime.datetime.now()
        text = (f"{Current_Time.year}/{Current_Time.month}/{Current_Time.day}\n"
                f"{self.Operator}")
        text_frame = shape.text_frame
        text_frame.clear()
        pg = text_frame.paragraphs[0]
        pg.font.size = Pt(32)
        pg.alignment = PP_ALIGN.CENTER
        pg.font.name = "メイリオ"
        # 色をグレーにする
        pg.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        pg.text = text

        print(f"\t>>> Creation day: {Current_Time.year}/{Current_Time.month}/{Current_Time.day}\n"
              f"\t>>> Operator: {self.Operator}")
        return None

    def add_slide(self):
        """
            受け取ったプレゼンテーションオブジェクトにスライドを追加し、追加されたスライドオブジェクトを返す。
            :return: スライドオブジェクト
            """
        # 白紙スライドの追加(ID=6は白紙スライド)
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        return slide

    def add_title(self, slide):
        """
        スライドに物質名を追加する
        :param slide: スライド
        :type slide: Slide object
        :return: なし
        """
        shapes = slide.shapes
        # MaterNameの文字数に応じて枠の横幅, 縦幅を調整する
        Box_Width = len(self.MaterName) * 230000
        Box_Height = 200000
        shape = shapes.add_textbox(50000, 50000, Box_Width, Box_Height)
        text_frame = shape.text_frame
        text_frame.clear()
        pg = text_frame.paragraphs[0]
        pg.font.size = Pt(10)
        pg.alignment = PP_ALIGN.CENTER
        pg.text = f"{self.MaterName}"
        pg.font.name = "メイリオ"
        # テキストボックスに枠線を追加
        line = shape.line
        # 枠線の色を黄緑にする
        line.color.rgb = RGBColor(0x00, 0x80, 0x00)
        line.width = Pt(1)
        # 文字を上下中央揃えにする
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        return

    @staticmethod
    def get_aspect(File_Name):
        """
        画像サイズを取得してアスペクト比を得る
        :param File_Name:　画像のファイル名
        :return: アスペクト比
        """
        im = Image.open(File_Name)
        im_width, im_height = im.size
        aspect_ratio = im_width / im_height
        return aspect_ratio

    def add_EvsAngles(self):
        """
        EvsAngles画像のファイル名を取得し、スライドに追加する

        IMG_DIRディレクトリから指定されたMaterNameに対応する
        EvsAngles画像を検索し、昇順にソートしてからスライドに
        貼り付けます。

        :return: なし
        :rtype: None
        """
        print("\n\tAdding Energy vs Angle images...")

        # EvsAnglesの画像のファイル名を取得
        EvsAngles = glob.glob(f"{self.IMG_DIR}/EvsAngle/{self.MaterName}_3molp*-12_EvsAngle.png")

        # pngで終了するファイル名のみ抽出。貼り付けたい画像の拡張子に応じて変える
        EvsAngles = [name for name in EvsAngles if name.endswith(".png")]

        # 昇順にソート（この順番でスライドに貼り付けられる）
        EvsAngles.sort()

        # スライドの作成
        self.mkEvsAngles(EvsAngles)

        print(f"\t>>> {Color.GREEN}Complete!{Color.RESET}")
        return

    def mkEvsAngles(self, Names):
        """
        指定された画像をスライドに追加します。

        画像リストから順にスライドに画像を追加し、配置を調整します。
        Debugモードが有効な場合、追加する画像の情報をコンソールに表示します。
        :param Names: 追加する画像ファイル名のリスト
        :type Names: list[str]
        :return: なし
        :rtype: None
        """
        # 新しいページを作成
        slide = self.add_slide()

        self.add_title(slide)

        # タイトルを追加
        shapes = slide.shapes
        shape = shapes.add_textbox(0, Constants.SLIDE_HEIGHT / 30, Constants.SLIDE_WIDTH, Constants.SLIDE_HEIGHT / 10)
        text_frame = shape.text_frame
        # 初期状態の空のパラグラフを使用してテキストを設定
        text_frame.clear()  # 空のパラグラフをクリア
        pg = text_frame.paragraphs[0]  # デフォルトのパラグラフを取得
        pg.font.size = Pt(36)
        pg.alignment = PP_ALIGN.CENTER
        pg.text = "Energy vs Angle"
        pg.font.name = "メイリオ"

        # コメントを追加
        shape = shapes.add_textbox(Constants.SLIDE_WIDTH * 0.05, Constants.IMG_CENTER_Y + Constants.IMG_CENTER_Y / 2,
                                   Constants.SLIDE_WIDTH * 0.9, Constants.SLIDE_HEIGHT / 5)
        text_frame = shape.text_frame
        text_frame.clear()
        pg = text_frame.paragraphs[0]
        pg.font.size = Pt(14)
        pg.font.name = "メイリオ"
        pg.text = "[Comment]"

        i = 0
        for name in Names:
            if self.debug:
                print(f"\t>>> Adding {name}...")
            aspect_ratio = self.get_aspect(name)
            # 画像のサイズを変換
            img_display_width = Constants.SLIDE_WIDTH / 3
            img_display_height = img_display_width / aspect_ratio

            # 画像の座標を計算
            left = img_display_width * i
            top = (Constants.SLIDE_HEIGHT - img_display_height) / 2 - Constants.SLIDE_HEIGHT / 15
            if self.debug:
                print(f"\t\t>>> left coordinate: {left}")
                print(f"\t\t>>> top coordinate: {top}")

            # 画像をスライドに追加
            slide.shapes.add_picture(name, left, top, width=img_display_width)
            i = i + 1
        return

    def add_TIs(self):
        """
        TIs画像のファイル名を取得し、スライドに追加する
        :return:
        """
        print("\n\tAdding Transfer Integral images...")

        # TIsの画像のファイル名を取得
        p1_TIs = glob.glob(f"{self.IMG_DIR}/TIvsAngle/{self.MaterName}_3molp1_{self.Tilt_Angle}-*_TIs.png")
        p2_TIs = glob.glob(f"{self.IMG_DIR}/TIvsAngle/{self.MaterName}_3molp2_{self.Tilt_Angle}-*_TIs.png")
        p3_TIs = glob.glob(f"{self.IMG_DIR}/TIvsAngle/{self.MaterName}_3molp3_{self.Tilt_Angle}-*_TIs.png")

        # pngで終了するファイル名のみ抽出。貼り付けたい画像の拡張子に応じて変える
        p1_TIs = [name for name in p1_TIs if name.endswith(".png")]
        p2_TIs = [name for name in p2_TIs if name.endswith(".png")]
        p3_TIs = [name for name in p3_TIs if name.endswith(".png")]

        # 昇順にソート（この順番でスライドに貼り付けられる）
        p1_TIs.sort()
        p2_TIs.sort()
        p3_TIs.sort()

        # スライドの作成
        self.mkTIs(p1_TIs, "p1")
        self.mkTIs(p2_TIs, "p2")
        self.mkTIs(p3_TIs, "p3")

        print(f"\t>>> {Color.GREEN}Complete!{Color.RESET}")
        return

    def mkTIs(self, Names, position):
        """
        指定された画像をスライドに追加し、位置とサイズを調整する

        この関数は、名前リストに基づいてスライドに画像を追加します。
        各画像のアスペクト比を計算し、スライド上で適切な位置と
        サイズに調整します。さらに、オプションでデバッグ情報を
        コンソールに出力します。
        :param list Names: 追加する画像のファイル名のリスト。
        :param str position: 画像の位置を示す文字列。

        :return: なし
        :rtype: None
        :raises FileNotFoundError: 指定された画像ファイルが存在しない場合に発生。
        """
        # 新しいページを作成
        slide = self.add_slide()
        self.add_title(slide)

        i = 0
        for name in Names:
            if self.debug:
                print(f"\t>>> Adding {name}...")
            aspect_ratio = self.get_aspect(name)

            # 画像のサイズを調整
            # スライドと画像のアスペクト比に応じて処理を分岐
            if aspect_ratio > Constants.SLIDE_ASPECT_RATIO:
                img_display_width = Constants.SLIDE_WIDTH / 2
                img_display_height = img_display_width / aspect_ratio
            else:  # 画像のほうが縦長だったら縦めいっぱいに広げる
                img_display_height = Constants.SLIDE_HEIGHT / 2
                img_display_width = img_display_height * aspect_ratio

            # 画像の座標を計算
            left, top = 0, 0
            # 12
            if i == 0:
                left = (Constants.IMG_CENTER_X - img_display_width) / 2
                top = (Constants.SLIDE_HEIGHT / 2 - img_display_height) / 2 + (Constants.SLIDE_HEIGHT / 15)
            # 23
            if i == 1:
                left = (Constants.IMG_CENTER_X - img_display_width) / 2
                top = (Constants.SLIDE_HEIGHT / 2 - img_display_height) / 2 + Constants.IMG_CENTER_Y
            # 31
            if i == 2:
                left = (Constants.IMG_CENTER_X - img_display_width) / 2 + Constants.IMG_CENTER_X
                top = (Constants.SLIDE_HEIGHT / 2 - img_display_height) / 2 + Constants.IMG_CENTER_Y

            if self.debug:
                print(f"\t\t>>> left coordinate: {left}")
                print(f"\t\t>>> top coordinate: {top}")

            # 画像をスライドに追加
            if aspect_ratio > Constants.SLIDE_ASPECT_RATIO:
                slide.shapes.add_picture(name, left, top, width=img_display_width)
            else:
                slide.shapes.add_picture(name, left, top, height=img_display_height)
            i = i + 1

        # タイトルを追加
        shapes = slide.shapes
        shape = shapes.add_textbox(0, Constants.SLIDE_HEIGHT / 30, Constants.SLIDE_WIDTH, Constants.SLIDE_HEIGHT / 10)
        text_frame = shape.text_frame
        # 初期状態の空のパラグラフを使用してテキストを設定
        text_frame.clear()  # 空のパラグラフをクリア
        pg = text_frame.paragraphs[0]  # デフォルトのパラグラフを取得
        pg.font.size = Pt(36)
        pg.alignment = PP_ALIGN.CENTER
        pg.text = f"Transfer Integral vs Angle - {position}"
        pg.font.name = "メイリオ"

        # コメントを追加
        shape = shapes.add_textbox(Constants.IMG_CENTER_X, Constants.SLIDE_HEIGHT / 7,
                                   Constants.IMG_CENTER_X, Constants.IMG_CENTER_Y - Constants.SLIDE_HEIGHT / 7)
        text_frame = shape.text_frame
        text_frame.clear()
        pg = text_frame.paragraphs[0]
        pg.font.size = Pt(14)
        pg.font.name = "メイリオ"
        pg.text = "[Comment]"
        return

    def add_EMvsAngles(self):
        """
        EMvsAnglesの画像ファイルを取得し、スライドに追加する。

        指定されたディレクトリ内のEMvsAnglesに関連する
        画像ファイル（HOMOおよびLUMO）を取得し、指定された
        スライドに追加します。画像ファイルはpng形式のみを対象とし、
        mkEMvsAngles関数を利用してスライドに配置します。

        :return: なし
        :rtype: None
        """
        print("\n\tAdding EMvsAngles images...")
        # EMvsAnglesの画像のファイル名を取得
        B12_EMvsAngles_HOMO = glob.glob(f"{self.IMG_DIR}/{self.MaterName}_{self.Tilt_Angle}-B12-HOMO_EMvsAngle.png")
        B12_EMvsAngles_LUMO = glob.glob(f"{self.IMG_DIR}/{self.MaterName}_{self.Tilt_Angle}-B12-LUMO_EMvsAngle.png")
        B3_EMvsAngles_HOMO = glob.glob(f"{self.IMG_DIR}/{self.MaterName}_{self.Tilt_Angle}-B3-HOMO_EMvsAngle.png")
        B3_EMvsAngles_LUMO = glob.glob(f"{self.IMG_DIR}/{self.MaterName}_{self.Tilt_Angle}-B3-LUMO_EMvsAngle.png")

        # pngで終了するファイル名のみ抽出。貼り付けたい画像の拡張子に応じて変える
        B12_EMvsAngles_HOMO = [name for name in B12_EMvsAngles_HOMO if name.endswith(".png")]
        B12_EMvsAngles_LUMO = [name for name in B12_EMvsAngles_LUMO if name.endswith(".png")]
        B3_EMvsAngles_HOMO = [name for name in B3_EMvsAngles_HOMO if name.endswith(".png")]
        B3_EMvsAngles_LUMO = [name for name in B3_EMvsAngles_LUMO if name.endswith(".png")]

        # スライドの作成
        self.mkEMvsAngles(B12_EMvsAngles_HOMO)
        self.mkEMvsAngles(B12_EMvsAngles_LUMO)
        self.mkEMvsAngles(B3_EMvsAngles_HOMO)
        self.mkEMvsAngles(B3_EMvsAngles_LUMO)

        print(f"\t>>> {Color.GREEN}Complete!{Color.RESET}")
        return

    def mkEMvsAngles(self, Names):
        """
        EMvsAngles画像をスライドに配置する。

        画像ファイルのリストに基づいて、新しいスライドを作成し、
        画像を配置します。スライドと画像のアスペクト比に応じて、
        画像のサイズを調整し、センタリングして追加します。

        :param Names: スライドに追加する画像ファイルのリスト
        :type Names: list[str]
        :return: なし
        :rtype: None
        """
        # 新しいページを作成
        slide = self.add_slide()
        self.add_title(slide)
        for name in Names:
            if self.debug:
                print(f"\t\tAdding {name}... ", end="")
            aspect_ratio = self.get_aspect(name)

            # タイトルを追加
            shapes = slide.shapes
            shape = shapes.add_textbox(0, Constants.SLIDE_HEIGHT / 30, Constants.SLIDE_WIDTH,
                                       Constants.SLIDE_HEIGHT / 10)
            text_frame = shape.text_frame
            # 初期状態の空のパラグラフを使用してテキストを設定
            text_frame.clear()  # 空のパラグラフをクリア
            pg = text_frame.paragraphs[0]  # デフォルトのパラグラフを取得
            pg.font.size = Pt(24)
            pg.alignment = PP_ALIGN.CENTER
            pg.text = f"Effective Mass, Energy vs Angle - {name.split('-')[1]} - {name.split('-')[2].split('_')[0]}"
            pg.font.name = "メイリオ"

            # スライドと画像のアスペクト比に応じて処理を分岐
            # 画像のほうが横長だったら横めいっぱいに広げる
            if aspect_ratio > Constants.SLIDE_ASPECT_RATIO:
                img_display_width = Constants.SLIDE_WIDTH
                img_display_height = img_display_width / aspect_ratio
            else:  # 画像のほうが縦長だったら縦めいっぱいに広げる
                img_display_height = Constants.SLIDE_HEIGHT
                img_display_width = img_display_height * aspect_ratio

            # センタリングする場合の画像の左上座標を計算
            left = Constants.IMG_CENTER_X - img_display_width / 2
            top = Constants.IMG_CENTER_Y - img_display_height / 2

            # 画像をスライドに追加
            if aspect_ratio > Constants.SLIDE_ASPECT_RATIO:
                slide.shapes.add_picture(name, left, top, width=img_display_width)
            else:
                slide.shapes.add_picture(name, left, top, height=img_display_height)
            if self.debug:
                print(f"{Color.GREEN}Complete!{Color.RESET}")
            return None

    def addStableBandFigures(self):
        """
        最安定構造におけるバンド図の画像をスライドに追加する。
        :return: なし
        """
        print("\n\tAdding Most Stable Band images...")

        # 最安定構造における有効質量を取得する
        with open(f"{self.MaterName}_{self.Tilt_Angle}_EffMasses.txt", "r") as f:
            # 最初の二行は読み飛ばす
            f.readline()
            f.readline()
            lines = f.readlines()
        Raw_Title = [line.split("\t")[0] for line in lines]
        # column方向の有効質量を取得する(from tensorの時は4、そうじゃない時は6に入っている)
        MCol_from_tensor = [line.split("\t")[4] for line in lines]
        MCol = [line.split("\t")[6] for line in lines]
        if self.debug:
            print(f"\t\t>>> Getting {Color.GREEN}Column{Color.RESET} Effective Masses...")
            print(f"\t\tTitle\tMCol_from_tensor\tMCol")
            for line in lines:
                split_line = line.split('\t')
                print(f"\t\t{split_line[0]}\t{split_line[4]}\t{split_line[6]}")
            print("")
        # MCol_from_tensorとMColを比較して、各行の要素が"-"ではない方をMColとして取得する
        MCol = [MCol[i] if MCol[i] != "-" else MCol_from_tensor[i] for i in range(len(MCol))]  # 後で調べる
        # Transverse方向の有効質量を取得する(from tensorの時は5、そうじゃない時は7に入っている)
        MTrv_from_tensor = [line.split("\t")[5] for line in lines]
        MTrv = [line.split("\t")[7] for line in lines]
        if self.debug:
            print(f"\t\t>>> Getting {Color.GREEN}Transverse{Color.RESET} Effective Masses...")
            print(f"\t\tTitle\tMC0l_from_tensor\tMCol")
            for line in lines:
                split_line = line.split('\t')
                print(f"\t\t{split_line[0]}\t{split_line[5]}\t{split_line[7]}")
            print("")
        # MTrv_from_tensorとMTrvを比較して、各行の要素が"-"ではない方をMTrvとして取得する
        MTrv = [MTrv[i] if MTrv[i] != "-" else MTrv_from_tensor[i] for i in range(len(MTrv))]  # 後で調べる
        if self.debug:
            print(f"\t************************* Effective Masses for B12 *************************")
            print("\tTitle\tMCol\tMTrv")
            for i in range(len(MCol)):
                if "-B12-" in Raw_Title[i]:
                    print(f"\t{Raw_Title[i]}\t{round(float(MCol[i]), 6)}\t{round(float(MTrv[i]), 6)}")
            print(f"\t****************************************************************************")

        # エネルギーを取得する
        # p1の最小エネルギーとその角度を取得する
        with open(f"{self.MaterName}_3molp1_{self.Tilt_Angle}_results/"
                  f"{self.MaterName}_3molp1_{self.Tilt_Angle}_min.txt", "r") as f:
            # 最初の二行は読み飛ばす
            f.readline()
            f.readline()
            lines = f.readlines()
        # linesの各行の要素の1番目(角度)を取得する
        angles = [line.split()[0] for line in lines]

        # linesの各行の４番目の要素(エネルギー)を取得する
        p1_band_energys = [line.split()[3] for line in lines]
        if self.debug:
            print(f"\n\t\t>>> Getting {Color.GREEN}p1{Color.RESET} Band Energies...")
            print(f"\t\tFile Name: {self.MaterName}_{self.Tilt_Angle}_3molp1_results/"
                  f"{self.MaterName}_{self.Tilt_Angle}_3molp1_min.txt")
            print(f"\t\tAngle\tEnergy")
            for angle in angles:
                print(f"\t\t{angle}\t{p1_band_energys[angles.index(angle)]}")

        # p2の最小エネルギーとその角度を取得する
        with open(f"{self.MaterName}_3molp2_{self.Tilt_Angle}_results/"
                  f"{self.MaterName}_3molp2_{self.Tilt_Angle}_min.txt", "r") as f:
            # 最初の二行は読み飛ばす
            f.readline()
            f.readline()
            lines = f.readlines()
        # linesの各行の要素の1番目(角度)を取得する
        angles = [line.split()[0] for line in lines]

        # linesの各行の４番目の要素(エネルギー)を取得する
        p2_band_energys = [line.split()[3] for line in lines]
        if self.debug:
            print(f"\n\t\t>>> Getting {Color.GREEN}p2{Color.RESET} Band Energies...")
            print(f"\t\tFile Name: {self.MaterName}_{self.Tilt_Angle}_3molp2_results/"
                  f"{self.MaterName}_{self.Tilt_Angle}_3molp2_min.txt")
            print(f"\t\tAngle\tEnergy")
            for angle in angles:
                print(f"\t\t{angle}\t{p2_band_energys[angles.index(angle)]}")

        # p1とp2の最小エネルギーを各角度ごとに足し合わせる
        p1_p2_band_energys = [float(p1_band_energys[i]) + float(p2_band_energys[i]) for i in
                              range(len(p1_band_energys))]
        if self.debug:
            print(f"\n"
                  f"\t\t>>> Getting {Color.GREEN}p1+p2{Color.RESET} Band Energies\n"
                  f"\t\tAngle\tP1 Energy\tP2 Energy\tP1+P2 Energy")
            for angle in angles:
                # 角度、p1のエネルギー、p2のエネルギー、p1+p2のエネルギーを表示
                print(f"\t\t{angle}\t{p1_band_energys[angles.index(angle)]}\t"
                      f"{p2_band_energys[angles.index(angle)]}\t{p1_p2_band_energys[angles.index(angle)]}")
        # p1_p2_band_energy'sの最小値とその要素が何番目の要素かを取得する
        min_p1_p2_band_energy = min(p1_p2_band_energys)
        min_index = p1_p2_band_energys.index(min_p1_p2_band_energy)
        # 最小値の角度を取得する
        min_angle = angles[min_index]
        if self.debug:
            print(f"\n"
                  f"\t\t{Color.GREEN}min_angle: {min_angle}{Color.RESET}")
            print(f"\t\t{Color.GREEN}min_band_figure: {min_p1_p2_band_energy}{Color.RESET}\n")

        # B12の最小エネルギーとその角度を取得する
        B12_HOMO_Min_MCol, B12_HOMO_Min_MTrv, B12_LUMO_Min_MCol, B12_LUMO_Min_MTrv, structure = 0, 0, 0, 0, []
        if self.debug:
            print(f"\t\tGetting Angle and Effective Mass for minimum Energy...")
        for i in range(len(Raw_Title)):
            parts = Raw_Title[i].split("/")[2].split("-")
            count = 0
            for part in parts:
                if part in ["B12", "B3"]:
                    structure = [part, count]
                count = count + 1
            Raw_Title[i].split("/")[2].split("-")[structure[1] + 1].replace("d", "")
            if structure[0] == "B12":
                Now_Angle = int(Raw_Title[i].split("/")[2].split("-")[structure[1] + 1].replace("d", ""))
                if self.debug:
                    print(f"\t\tNow_Angle: {Now_Angle}\tmin_angle: {min_angle}\t{Now_Angle == int(float(min_angle))}")
                if Now_Angle == int(float(min_angle)):
                    if self.debug:
                        print(f'\t\t{Raw_Title[i].split("/")[2].split("-")[3].split(".")[0]}')
                    if Raw_Title[i].split("/")[2].split("-")[structure[1] + 2].split(".")[0] == "HOMO":
                        B12_HOMO_Min_MCol = MCol[i]
                        B12_HOMO_Min_MTrv = MTrv[i]
                        if self.debug:
                            print(f"\t\tMCol: {B12_HOMO_Min_MCol}\t\tMTrv: {B12_HOMO_Min_MTrv}\t\t"
                                  f"Mcol/Mtrv: {float(B12_HOMO_Min_MCol) / float(B12_HOMO_Min_MTrv)}\n")
                    elif Raw_Title[i].split("/")[2].split("-")[structure[1] + 2].split(".")[0] == "LUMO":
                        B12_LUMO_Min_MCol = MCol[i]
                        B12_LUMO_Min_MTrv = MTrv[i]
                        if self.debug:
                            print(f"\t\tMCol: {B12_LUMO_Min_MCol}\t\tMTrv: {B12_LUMO_Min_MTrv}\t\t"
                                  f"Mcol/Mtrv: {float(B12_LUMO_Min_MCol) / float(B12_LUMO_Min_MTrv)}\n")

            else:
                pass
        self.mkStableBandFigures(12, min_angle, min_p1_p2_band_energy,
                                 B12_HOMO_Min_MCol, B12_HOMO_Min_MTrv, B12_LUMO_Min_MCol, B12_LUMO_Min_MTrv)

        if self.debug:
            print(f"\t************************* Effective Masses for B3 *************************")
            print("\tTitle\tMCol\tMTrv")
            for i in range(len(MCol)):
                if "-B3-" in Raw_Title[i]:
                    print(f"\t{Raw_Title[i]}\t{round(float(MCol[i]), 6)}\t{round(float(MTrv[i]), 6)}")
            print(f"\t****************************************************************************")

        # p3の最小エネルギーとその角度を取得する
        with open(f"{self.MaterName}_3molp3_{self.Tilt_Angle}_results/"
                  f"{self.MaterName}_3molp3_{self.Tilt_Angle}_min.txt", "r") as f:
            # 最初の二行は読み飛ばす
            f.readline()
            f.readline()
            lines = f.readlines()
        # linesの各行の要素の1番目(角度)を取得する
        angles = [line.split()[0] for line in lines]

        # linesの各行の４番目の要素(エネルギー)を取得する
        band_energy = [line.split()[3] for line in lines]
        if self.debug:
            print(f"\n\t\t>>> Getting {Color.GREEN}p3{Color.RESET} Band Energies...")
            print(f"\t\tFile Name: {self.MaterName}_{self.Tilt_Angle}_3molp3_results/"
                  f"{self.MaterName}_{self.Tilt_Angle}_3molp3_min.txt")
            print(f"\t\tAngle\tEnergy")
            for angle in angles:
                print(f"\t\t{angle}\t{band_energy[angles.index(angle)]}")

        # band_figuresの最小値とその要素が何番目の要素かを取得する
        min_band_energy = min(band_energy)
        min_index = band_energy.index(min_band_energy)
        # 最小値の角度を取得する
        min_angle = angles[min_index]
        if self.debug:
            print(f"\n"
                  f"\t\t{Color.GREEN}min_angle: {min_angle}{Color.RESET}")
            print(f"\t\t{Color.GREEN}min_band_figure: {min_p1_p2_band_energy}{Color.RESET}\n")

        B3_HOMO_Min_MCol, B3_HOMO_Min_MTrv, B3_LUMO_Min_MCol, B3_LUMO_Min_MTrv = 0, 0, 0, 0
        for i in range(len(Raw_Title)):
            parts = Raw_Title[i].split("/")[2].split("-")
            count = 0
            for part in parts:
                if part in ["B12", "B3"]:
                    structure = [part, count]
                count = count + 1
            Raw_Title[i].split("/")[2].split("-")[structure[1] + 1].replace("d", "")
            if structure[0] == "B3":
                Now_Angle = int(Raw_Title[i].split("/")[2].split("-")[structure[1] + 1].replace("d", ""))
                if self.debug:
                    print(f"\t\tNow_Angle: {Now_Angle}\tmin_angle: {min_angle}\t{Now_Angle == int(float(min_angle))}")
                if Now_Angle == int(float(min_angle)):
                    if self.debug:
                        print(f"\t\t{Raw_Title[i].split("/")[2].split("-")[3].split(".")[0]}")
                    if Raw_Title[i].split("/")[2].split("-")[structure[1] + 2].split(".")[0] == "HOMO":
                        B3_HOMO_Min_MCol = MCol[i]
                        B3_HOMO_Min_MTrv = MTrv[i]
                        if self.debug:
                            print(f"\t\tMCol: {B3_HOMO_Min_MCol}\t\tMTrv: {B3_HOMO_Min_MTrv}\t\t"
                                  f"Mcol/Mtrv: {float(B3_HOMO_Min_MCol) / float(B3_HOMO_Min_MTrv)}\n")
                    elif Raw_Title[i].split("/")[2].split("-")[structure[1] + 2].split(".")[0] == "LUMO":
                        B3_LUMO_Min_MCol = MCol[i]
                        B3_LUMO_Min_MTrv = MTrv[i]
                        if self.debug:
                            print(f"\t\tMCol: {B3_LUMO_Min_MCol}\t\tMTrv: {B3_LUMO_Min_MTrv}\t\t"
                                  f"Mcol/Mtrv: {float(B3_LUMO_Min_MCol) / float(B3_LUMO_Min_MTrv)}\n")

        # 最安定構造におけるバンド図の画像をスライドに追加する
        self.mkStableBandFigures(3, min_angle, float(min_band_energy) * 2,
                                 B3_HOMO_Min_MCol, B3_HOMO_Min_MTrv, B3_LUMO_Min_MCol, B3_LUMO_Min_MTrv)

        print(f"\t>>> {Color.GREEN}Complete!{Color.RESET}")
        return None

    def mkStableBandFigures(self, position, Angle, min_band_energy,
                            HOMO_MCol, HOMO_MTrv, LUMO_MCol, LUMO_MTrv):
        """
        最安定構造におけるバンド図の画像をスライドに追加する。
        :param position: 1,2,3のいずれか
        :param Angle: 角度
        :param min_band_energy: 最小のエネルギー
        :param LUMO_MTrv:

        :param LUMO_MCol:
        :param HOMO_MTrv:
        :param HOMO_MCol:
        :return: なし
        """
        if self.debug:
            print(f"\n\t\t>>> Making Stable Band Figures...")

        # 新しいページを作成
        slide = self.add_slide()

        self.add_title(slide)

        # もしpositionが1,2ならB12構造、それ以外ならB3構造
        if position == 1 or position == 2:
            Structure = "B12"
        else:
            Structure = "B3"
        H_File_Name = (f"Figures/BandStructures/{self.MaterName}_{self.Tilt_Angle}"
                       f"-{Structure}-{int(float(Angle))}d-H-2.png")
        # H_File_Nameが存在するかを確認
        try:
            with open(H_File_Name):
                pass
        except FileNotFoundError:
            print(f"{H_File_Name} is not found.")
            return
        L_File_Name = (f"Figures/BandStructures/{self.MaterName}_{self.Tilt_Angle}"
                       f"-{Structure}-{int(float(Angle))}d-L-2.png")
        # L_File_Nameが存在するかを確認
        try:
            with open(L_File_Name):
                pass
        except FileNotFoundError:
            print(f"{L_File_Name} is not found.")
            return
        if self.debug:
            print(f"\t\t>>> HOMO Picture File Name: {H_File_Name}")
            print(f"\t\t>>> LUMO Picture File Name: {L_File_Name}")

        # 画像のサイズを取得
        aspect_ratio_H = self.get_aspect(H_File_Name)
        aspect_ratio_L = self.get_aspect(L_File_Name)
        # スライドと画像のアスペクト比に応じて処理を分岐
        # 画像のほうが横長だったら横めいっぱいに広げる
        if aspect_ratio_H > Constants.SLIDE_ASPECT_RATIO:
            img_display_width_H = Constants.SLIDE_WIDTH / 2
            img_display_height_H = img_display_width_H / aspect_ratio_H
        else:  # 画像のほうが縦長だったら縦めいっぱいに広げる
            img_display_height_H = Constants.SLIDE_HEIGHT / 2
            img_display_width_H = img_display_height_H * aspect_ratio_H
        if aspect_ratio_L > Constants.SLIDE_ASPECT_RATIO:
            img_display_width_L = Constants.SLIDE_WIDTH / 2
            img_display_height_L = img_display_width_L / aspect_ratio_L
        else:  # 画像のほうが縦長だったら縦めいっぱいに広げる
            img_display_height_L = Constants.SLIDE_HEIGHT / 2
            img_display_width_L = img_display_height_L * aspect_ratio_L

        # センタリングする場合の画像の左上座標を計算
        left_H = Constants.IMG_CENTER_X - img_display_width_H
        top_H = Constants.IMG_CENTER_Y - img_display_height_H / 2 + Constants.SLIDE_HEIGHT / 18
        left_L = Constants.IMG_CENTER_X
        top_L = Constants.IMG_CENTER_Y - img_display_height_L / 2 + Constants.SLIDE_HEIGHT / 18

        # 画像をスライドに追加
        slide.shapes.add_picture(H_File_Name, left_H, top_H, width=img_display_width_H)
        slide.shapes.add_picture(L_File_Name, left_L, top_L, width=img_display_width_L)

        # タイトルを追加
        shapes = slide.shapes
        shape = shapes.add_textbox(0, Constants.SLIDE_HEIGHT / 30, Constants.SLIDE_WIDTH, Constants.SLIDE_HEIGHT / 10)
        text_frame = shape.text_frame
        # 初期状態の空のパラグラフを使用してテキストを設定
        text_frame.clear()  # 空のパラグラフをクリア
        pg = text_frame.paragraphs[0]  # デフォルトのパラグラフを取得
        pg.font.size = Pt(36)
        pg.alignment = PP_ALIGN.CENTER
        pg.text = f"エネルギーが最小のバンド構造 - B{position}"
        pg.font.name = "メイリオ"

        # タイトルの下に角度、エネルギー、MCol、MTrvを追加
        shape = shapes.add_textbox(Constants.SLIDE_WIDTH * 0.025, Constants.SLIDE_HEIGHT / 8,
                                   Constants.SLIDE_WIDTH * 0.95, Constants.SLIDE_HEIGHT / 8)
        text_frame = shape.text_frame
        text_frame.clear()
        pg = text_frame.paragraphs[0]
        pg.font.size = Pt(14)
        pg.font.name = "メイリオ"
        pg.text = (f"Angle: {Angle}\n"
                   f"Energy: {min_band_energy}\n"
                   f"HOMO Mcol\t\tHOMO Mtrv\t\tMcol/Mtrv\n"
                   f"{round(float(HOMO_MCol), 3)}\t\t{round(float(HOMO_MTrv), 3)}\t\t"
                   f"{round(float(HOMO_MCol) / float(HOMO_MTrv), 3)}\n"
                   f"LUMO Mcol\t\tLUMO Mtrv\t\tMcol/Mtrv\n"
                   f"{round(float(LUMO_MCol), 3)}\t\t{round(float(LUMO_MTrv), 3)}\t\t"
                   f"{round(float(LUMO_MCol) / float(LUMO_MTrv), 3)}")

        # 画像の下にコメントを追加
        shape = shapes.add_textbox(Constants.SLIDE_WIDTH * 0.1,
                                   Constants.SLIDE_HEIGHT / 2 + img_display_height_H / 2 + Constants.SLIDE_HEIGHT / 18,
                                   Constants.SLIDE_WIDTH * 0.8, Constants.SLIDE_HEIGHT / 8)
        text_frame = shape.text_frame
        text_frame.clear()
        pg = text_frame.paragraphs[0]
        pg.font.size = Pt(14)
        pg.font.name = "メイリオ"
        pg.text = "[Comment]"

        return


class Constants:
    # ハートリーからeVに変換するための定数
    Hartree_to_eV = 27.21138602
    # 気体定数[R] (J/(mol・K))
    R = 8.314462618
    # 真空中の光速[m/s]
    C = 299792458
    # 電気素量[C]
    E = 1.602176634e-19
    # meV/(m/s)^2
    ElMass = 0.51099895 * 10 ** 9 / C ** 2
    # meVs
    h_bar = 6.582119569 * 10 ** (-13)
    # 近似のための分割数
    n = 100
    # プロットパターン
    Pattern = "2"

    # スライドサイズ
    # 4:3 (default) 9144000x6858000
    # 16:9 12193200x6858000
    SLIDE_WIDTH, SLIDE_HEIGHT = 9144000, 6858000
    # スライド中心のX、Y座標（左上が原点）
    IMG_CENTER_X, IMG_CENTER_Y = SLIDE_WIDTH / 2, SLIDE_HEIGHT / 2
    # スライドのアスペクト比
    SLIDE_ASPECT_RATIO = SLIDE_WIDTH / SLIDE_HEIGHT


class Color:
    """
    ANSI escape code for color
    """

    def __init__(self):
        self._Color = "Color"

    BLACK = '\033[30m'  # black
    RED = '\033[31m'  # red
    GREEN = '\033[32m'  # green
    YELLOW = '\033[33m'  # yellow
    BLUE = '\033[34m'  # blue
    MAGENTA = '\033[35m'  # magenta
    CYAN = '\033[36m'  # cyan
    LIME = '\033[92m'  # lime
    WHITE = '\033[37m'  # white
    COLOR_DEFAULT = '\033[39m'  # default
    BOLD = '\033[1m'  # bold
    UNDERLINE = '\033[4m'  # underline
    INVISIBLE = '\033[08m'  # invisible
    REVERSE = '\033[07m'  # reverse
    BG_BLACK = '\033[40m'  # (background)black
    BG_RED = '\033[41m'  # (background)red
    BG_GREEN = '\033[42m'  # (background)green
    BG_YELLOW = '\033[43m'  # (background)yellow
    BG_BLUE = '\033[44m'  # (background)blue
    BG_MAGENTA = '\033[45m'  # (background)magenta
    BG_CYAN = '\033[46m'  # (background)cyan
    BG_WHITE = '\033[47m'  # (background)white
    BG_DEFAULT = '\033[49m'  # (background)default
    RESET = '\033[0m'  # reset


class StandardPhrases:
    def __init__(self):
        self._StandardPhrases = "StandardPhrases"

    ProgramAbst = ("\n"
                   "*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*\n"
                   "*   This program collects and processes calculation results of                *\n"
                   "*       optimized aggregate structures of organic semiconductors              *\n"
                   "*       with HerringBone (HB) structures.                                     *\n"
                   "*   It performs the following tasks:                                          *\n"
                   "*       - Calculates effective masses                                         *\n"
                   "*       - Plots band diagrams                                                 *\n"
                   "*       - Generates summary slides of the results                             *\n"
                   "*                                                                             *\n"
                   "*   To run this program, please ensure the following files are available:     *\n"
                   "*       - “MoleculeName”_3molp1_t”N”d_results                                 *\n"
                   "*       - “MoleculeName”_3molp2_t”N”d_results                                 *\n"
                   "*       - “MoleculeName”_3molp3_t”N”d_results                                 *\n"
                   "*                                                                             *\n"
                   "*   Replace “MoleculeName” with the actual name of your molecule,             *\n"
                   "*       and “N” with the appropriate value used in your calculations.         *\n"
                   "*                                                                             *\n"
                   "*   When you start the program, you will be prompted to enter your name,      *\n"
                   "*       - which will be used as the slide creator’s name.                     *\n"
                   "*                                                                             *\n"
                   "*                                 created by Toshiyuki Togashi 2024/11/27     *\n"
                   "*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*-*\n")
    AbnormalEnd = "\n************* Programme DID NOT terminate successfully. *************\n"
    HelpText = ("\n"
                "The required file may not exist.\n"
                "Please check.")


if __name__ == '__main__':
    main()
